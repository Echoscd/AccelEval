#!/usr/bin/env python3
"""Generate ORBench group-meeting presentation for GPU-acceleration newcomers."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── colour palette ──────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
BG_MID    = RGBColor(0x22, 0x22, 0x3A)
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)    # cyan
ACCENT2   = RGBColor(0xFF, 0x6B, 0x6B)    # coral
ACCENT3   = RGBColor(0x4E, 0xCB, 0x71)    # green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xDD)
GRAY      = RGBColor(0x88, 0x88, 0x99)
ORANGE    = RGBColor(0xFF, 0xA5, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── helpers ─────────────────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6),
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p

def add_bullet_slide(slide, items, left=Inches(0.8), top=Inches(2.0),
                     width=Inches(11.5), size=20, color=LIGHT):
    """Add a list of bullet items. Each item is (text,) or (text, sub_items)."""
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            main_text, subs = item[0], item[1] if len(item) > 1 else []
        else:
            main_text, subs = item, []
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"  {main_text}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        for sub in subs:
            p2 = tf.add_paragraph()
            p2.text = f"      {sub}"
            p2.font.size = Pt(size - 3)
            p2.font.color.rgb = GRAY
            p2.font.name = "Calibri"
            p2.space_before = Pt(2)

def title_bar(slide, title, subtitle=None):
    add_shape(slide, 0, 0, W, Inches(1.4), BG_MID)
    add_shape(slide, 0, Inches(1.38), W, Inches(0.04), ACCENT)
    add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             title, size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.45),
                 subtitle, size=16, color=GRAY)

# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s, BG_DARK)
add_shape(s, 0, Inches(2.6), W, Inches(2.8), BG_MID)
add_text(s, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
         "ORBench", size=60, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "LLM 能写出多快的 GPU 代码？", size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
         "通用 CPU -> CUDA 加速 Benchmark：评测大语言模型的 GPU 编程能力",
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
         "组会报告  |  2026.04", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2 — Outline
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "Outline")
items = [
    "1.  背景：为什么需要 GPU？为什么让 LLM 来写？",
    "2.  现有 Benchmark 的不足",
    "3.  ORBench 设计理念与架构",
    "4.  任务总览：40 道题，10+ 领域",
    "5.  三层评测框架：如何公平比赛",
    "6.  Prompt 分级与多轮 Agent",
    "7.  初步实验结果与发现",
    "8.  总结与下一步",
]
add_bullet_slide(s, items, size=22, color=LIGHT)

# =====================================================================
# SLIDE 3 — Why GPU?
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "为什么需要 GPU？", "从一个简单的类比说起")

# CPU vs GPU analogy
add_shape(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.2), BG_MID)
add_text(s, Inches(1.0), Inches(1.9), Inches(5), Inches(0.5),
         "CPU — 数学教授", size=24, color=ACCENT, bold=True)
tf = add_text(s, Inches(1.0), Inches(2.5), Inches(5), Inches(1.4),
              "6~64 个核心，每个都很聪明", size=17, color=LIGHT)
add_para(tf, "擅长复杂逻辑、分支判断、串行任务", size=17, color=LIGHT)
add_para(tf, "一个人算微积分 —— 快，但只有一个人", size=17, color=LIGHT)

add_shape(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.2), BG_MID)
add_text(s, Inches(7.2), Inches(1.9), Inches(5), Inches(0.5),
         "GPU — 10,000 名小学生", size=24, color=ACCENT2, bold=True)
tf = add_text(s, Inches(7.2), Inches(2.5), Inches(5), Inches(1.4),
              "数千个核心，每个都很简单", size=17, color=LIGHT)
add_para(tf, "擅长大规模重复计算、数据并行", size=17, color=LIGHT)
add_para(tf, "一万人同时算加法 —— 每人慢，但人多", size=17, color=LIGHT)

# Key insight
add_shape(s, Inches(2.5), Inches(4.4), Inches(8.3), Inches(1.2), BG_MID)
add_text(s, Inches(2.7), Inches(4.5), Inches(7.8), Inches(0.5),
         "关键洞察", size=22, color=ACCENT3, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.7), Inches(5.0), Inches(7.8), Inches(0.5),
         "很多科学计算和工程问题天然适合「让一万个人同时算」", size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Numbers
add_text(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5),
         "NVIDIA H100:  16,896 CUDA cores  |  3.35 TFLOPS (FP64)  |  vs CPU ~0.05 TFLOPS",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
         "但实际加速比取决于算法的并行度、内存访问模式、同步开销 —— 这正是 ORBench 要评测的",
         size=14, color=ORANGE, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 4 — CUDA 编程入门
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "CUDA 编程：30 秒入门", "把 for 循环变成并行")

# CPU code
add_shape(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.8), BG_MID)
add_text(s, Inches(0.7), Inches(1.85), Inches(5.5), Inches(0.4),
         "CPU 串行代码", size=18, color=ACCENT2, bold=True)
code_cpu = (
    "// CPU: 逐个计算\n"
    "for (int i = 0; i < N; i++) {\n"
    "    C[i] = A[i] + B[i];\n"
    "}\n"
    "// N=1000万 -> 循环 1000万次"
)
add_text(s, Inches(0.7), Inches(2.35), Inches(5.5), Inches(2.0),
         code_cpu, size=15, color=LIGHT, font_name="Consolas")

# GPU code
add_shape(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.8), BG_MID)
add_text(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(0.4),
         "GPU 并行代码 (CUDA)", size=18, color=ACCENT3, bold=True)
code_gpu = (
    "// GPU: 每个线程算一个\n"
    "__global__ void add(float *A,\n"
    "                    float *B, float *C) {\n"
    "    int i = threadIdx.x + blockIdx.x\n"
    "            * blockDim.x;\n"
    "    C[i] = A[i] + B[i];\n"
    "}\n"
    "// N=1000万 -> 1000万线程同时算"
)
add_text(s, Inches(7.2), Inches(2.35), Inches(5.5), Inches(2.2),
         code_gpu, size=15, color=LIGHT, font_name="Consolas")

# Challenge box
add_shape(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), BG_MID)
add_text(s, Inches(0.7), Inches(5.1), Inches(11.8), Inches(0.4),
         "写 CUDA 的真正难点不是语法，而是...", size=20, color=ORANGE, bold=True)
tf = add_text(s, Inches(0.7), Inches(5.6), Inches(11.8), Inches(1.2),
              "  内存合并访问 (Coalesced Access) —— 相邻线程访问相邻内存才快", size=16, color=LIGHT)
add_para(tf, "  共享内存与 Bank Conflict —— GPU 的 L1 cache，用好了快 10x", size=16, color=LIGHT)
add_para(tf, "  线程同步与原子操作 —— 一万个人同时写同一个变量会出问题", size=16, color=LIGHT)
add_para(tf, "  Warp Divergence —— 32 个线程必须执行相同指令，if-else 会让一半人等另一半", size=16, color=LIGHT)

# =====================================================================
# SLIDE 5 — Why LLM + GPU?
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "为什么让 LLM 来写 GPU 代码？", "一个真实的痛点")

items = [
    ("手写 CUDA 的现状", [
        "CUDA 专家稀缺，培养周期长（6-12 个月）",
        "同一个算法，不同硬件（A100 vs H100）最优实现不同",
        "调优一个 kernel 可能需要数天，涉及大量 trial-and-error",
    ]),
    ("LLM 的潜力", [
        "GPT-4、Claude、Gemini 已经能写出语法正确的 CUDA 代码",
        "关键问题：语法正确 != 高效。LLM 写的 kernel 到底有多快？",
        "有没有可能让 LLM 自动迭代优化，像人类工程师一样调优？",
    ]),
    ("学术价值：这是一个未被充分评测的能力维度", [
        "代码生成 Benchmark 很多（HumanEval、MBPP），但几乎不涉及 GPU",
        "GPU Benchmark 很少，且覆盖面窄 → 这正是 ORBench 要解决的",
    ]),
]
add_bullet_slide(s, items, size=19, color=LIGHT)

# =====================================================================
# SLIDE 6 — Existing benchmarks
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "现有 GPU Benchmark 的不足")

# Table header
y0 = Inches(1.8)
row_h = Inches(0.5)
col_w = [Inches(2.8), Inches(4.0), Inches(5.3)]
col_x = [Inches(0.5)]
for i in range(1, 3):
    col_x.append(col_x[-1] + col_w[i-1] + Inches(0.1))

headers = ["Benchmark", "覆盖领域", "局限"]
for j, h in enumerate(headers):
    add_shape(s, col_x[j], y0, col_w[j], row_h, ACCENT)
    add_text(s, col_x[j], y0, col_w[j], row_h, f"  {h}", size=16, color=BG_DARK, bold=True)

rows = [
    ["KernelBench", "ML 算子 (MatMul, Conv, Softmax)", "只覆盖 ML，不评测通用算法"],
    ["ComputeEval", "CUDA 语法练习题", "只检查编译/正确性，不评性能"],
    ["HumanEval-X", "简单编程题 (多语言)", "无 GPU 题目，无性能指标"],
    ["ORBench (Ours)", "图算法/DP/金融/几何/科学计算...", "40 题 10+ 领域, 评正确性+性能+加速比"],
]
for i, row in enumerate(rows):
    y = y0 + row_h * (i + 1) + Inches(0.05) * (i + 1)
    bg = BG_MID if i < 3 else RGBColor(0x1E, 0x3A, 0x2E)
    tc = LIGHT if i < 3 else ACCENT3
    for j, cell in enumerate(row):
        add_shape(s, col_x[j], y, col_w[j], row_h, bg)
        add_text(s, col_x[j], y, col_w[j], row_h,
                 f"  {cell}", size=14, color=tc, bold=(i == 3))

# =====================================================================
# SLIDE 7 — ORBench Design Philosophy
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "ORBench 设计理念", "公平、全面、自动化")

boxes = [
    ("覆盖面广", "40 道任务\n10+ 算法领域\n难度 1-4 星", ACCENT),
    ("公平评测", "统一 harness\n防作弊机制\nCUDA Event 精确计时", ACCENT2),
    ("自动化流水线", "生成 -> 编译 -> 验证\n-> 基准测试 -> 分析\n一条命令搞定", ACCENT3),
    ("多模型对比", "Claude / GPT-4o\nGemini / DeepSeek\n同题同评", ORANGE),
]
for i, (title, desc, color) in enumerate(boxes):
    x = Inches(0.5 + i * 3.2)
    add_shape(s, x, Inches(2.0), Inches(2.8), Inches(3.5), BG_MID)
    add_shape(s, x, Inches(2.0), Inches(2.8), Inches(0.06), color)
    add_text(s, x, Inches(2.3), Inches(2.8), Inches(0.5),
             title, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.0), Inches(2.8), Inches(2.0),
             desc, size=16, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
         "核心问题：LLM 不仅要写对，还要写快 —— 我们同时评测正确性和性能",
         size=18, color=ORANGE, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 8 — Three-layer architecture
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "三层评测架构", "LLM 只需要写纯计算逻辑，I/O 和计时由框架处理")

layers = [
    ("Harness 层 (harness_common.h)", "框架提供", "计时、Warmup、验证控制\nCUDA Event / clock_gettime 精确计时", ACCENT),
    ("Task I/O 层 (task_io.cu)", "任务特定", "解析 input.bin、桥接 harness 与 solution\n格式化输出、对比 expected_output", ACCENT2),
    ("Solution 层 (LLM 生成)", "LLM 输出", "纯算法实现：solution_init() + solution_compute()\n不需要处理文件 I/O、计时、验证", ACCENT3),
]
for i, (name, tag, desc, color) in enumerate(layers):
    y = Inches(1.8 + i * 1.7)
    add_shape(s, Inches(0.8), y, Inches(11.5), Inches(1.4), BG_MID)
    add_shape(s, Inches(0.8), y, Inches(0.06), Inches(1.4), color)
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(6), Inches(0.4),
             name, size=20, color=color, bold=True)
    add_text(s, Inches(1.2), y + Inches(0.55), Inches(8), Inches(0.8),
             desc, size=15, color=LIGHT)
    # tag
    add_shape(s, Inches(10.0), y + Inches(0.3), Inches(2.0), Inches(0.5), color)
    add_text(s, Inches(10.0), y + Inches(0.3), Inches(2.0), Inches(0.5),
             tag, size=14, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

# Arrow hints
add_text(s, Inches(5.5), Inches(3.15), Inches(2), Inches(0.3),
         "         v  调用  v", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(5.5), Inches(4.85), Inches(2), Inches(0.3),
         "         v  调用  v", size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
         "好处：LLM 不需要写 boilerplate，我们也不用担心它通过 I/O 作弊", size=16, color=ORANGE)

# =====================================================================
# SLIDE 9 — Task Overview (Categories)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "任务总览：40 道题，10+ 领域")

categories = [
    ("金融计算 (8)", "Black-Scholes, 债券定价\n蒙特卡洛, 亚式期权\n动态定价 DP ...", ACCENT),
    ("图算法 (6)", "Bellman-Ford, Floyd-Warshall\nPageRank, 最大流\n交通分配 ...", ACCENT2),
    ("动态规划 (8)", "收益管理 DP, TSP\nViterbi HMM, CYK 解析\n库存优化 ...", ACCENT3),
    ("科学计算 (6)", "N-Body, SPH 流体\nGROMACS 分子力\n热扩散 Jacobi ...", ORANGE),
    ("空间/序列 (5)", "DTW 时间弯曲\nSmith-Waterman\nDBSCAN 聚类 ...", RGBColor(0xBB, 0x86, 0xFC)),
    ("其他 (7)", "碰撞检测, SpMV\nLP 求解器, 正则匹配\n组合优化 ...", RGBColor(0xFF, 0xD5, 0x4F)),
]
for i, (cat, examples, color) in enumerate(categories):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.8 + row * 2.6)
    add_shape(s, x, y, Inches(3.8), Inches(2.2), BG_MID)
    add_shape(s, x, y, Inches(3.8), Inches(0.05), color)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4),
             cat, size=18, color=color, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.65), Inches(3.4), Inches(1.4),
             examples, size=14, color=LIGHT)

# =====================================================================
# SLIDE 10 — Difficulty distribution
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "难度分布与并行模式")

# Difficulty bars (horizontal)
diffs = [
    ("*    简单 (9题)", 9, "Black-Scholes, SpMV, 欧氏距离矩阵 ...", ACCENT3),
    ("**   中等 (10题)", 10, "Bellman-Ford, PageRank, 热扩散, 蒙特卡洛 ...", ACCENT),
    ("***  困难 (14题)", 14, "碰撞检测, DTW, N-Body, DBSCAN, SPH ...", ORANGE),
    ("**** 极难 (7题)", 7, "收益管理 DP, TSP, CYK 解析, SPH 力计算 ...", ACCENT2),
]
for i, (label, count, examples, color) in enumerate(diffs):
    y = Inches(1.8 + i * 1.2)
    add_text(s, Inches(0.8), y, Inches(2.5), Inches(0.4),
             label, size=18, color=color, bold=True)
    bar_w = Inches(count * 0.45)
    add_shape(s, Inches(3.5), y + Inches(0.05), bar_w, Inches(0.3), color)
    add_text(s, Inches(3.5) + bar_w + Inches(0.2), y, Inches(0.5), Inches(0.4),
             str(count), size=16, color=color, bold=True)
    add_text(s, Inches(3.5), y + Inches(0.45), Inches(9), Inches(0.35),
             examples, size=13, color=GRAY)

# Parallelism patterns
add_shape(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), BG_MID)
add_text(s, Inches(0.7), Inches(6.05), Inches(11.8), Inches(0.4),
         "覆盖的并行模式：数据并行 / Wavefront / 规约 / Scan / Tiling / 原子操作 / 持久 Kernel",
         size=16, color=LIGHT)
add_text(s, Inches(0.7), Inches(6.45), Inches(11.8), Inches(0.4),
         "输入规模：10K ~ 10M 数据点 —— 足够让 GPU 「吃饱」",
         size=16, color=LIGHT)

# =====================================================================
# SLIDE 11 — Interface modes
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "双接口模式：防作弊设计", "不让 LLM 把计算藏在初始化里")

# Two mode boxes
for j, (mode, when, timing, trick, color) in enumerate([
    ("init_compute", "输入数据大 (>1MB)\n如大规模图", "只计时 compute()",
     "允许 init 预处理\n（建索引、预排序）", ACCENT),
    ("compute_only", "输入数据小 (<1MB)\n如矩阵运算", "计时 setup + compute()",
     "防止把计算逻辑\n藏在 init 里作弊", ACCENT2),
]):
    x = Inches(0.8 + j * 6.3)
    add_shape(s, x, Inches(1.8), Inches(5.8), Inches(3.8), BG_MID)
    add_shape(s, x, Inches(1.8), Inches(5.8), Inches(0.06), color)
    add_text(s, x + Inches(0.3), Inches(2.0), Inches(5), Inches(0.5),
             mode, size=26, color=color, bold=True, font_name="Consolas")

    labels = ["适用场景", "计时范围", "特点"]
    values = [when, timing, trick]
    for k, (lbl, val) in enumerate(zip(labels, values)):
        yy = Inches(2.7 + k * 0.9)
        add_text(s, x + Inches(0.3), yy, Inches(1.8), Inches(0.3),
                 lbl, size=14, color=GRAY)
        add_text(s, x + Inches(2.2), yy, Inches(3.2), Inches(0.8),
                 val, size=15, color=LIGHT)

add_text(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
         "通过 task.json 的 interface_mode 字段控制，编译时自动传 -DORBENCH_COMPUTE_ONLY 宏",
         size=14, color=GRAY)

# =====================================================================
# SLIDE 12 — Prompt Levels
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "Prompt 三级分级", "从手把手教到完全放手")

levels = [
    ("L1 — 保姆级", "任务描述 + 接口定义\n+ 算法背景 + CPU 代码\n+ 详细 GPU 优化提示",
     "评测：代码实现能力\n「告诉你怎么优化，你能写出来吗？」", ACCENT3),
    ("L2 — 适中", "任务描述 + 接口定义\n+ CPU 代码 + 简要提示",
     "评测：优化策略选择\n「给个方向，你能找到最佳方案吗？」", ACCENT),
    ("L3 — 放手", "任务描述 + 接口定义\n+ CPU 代码（无提示）",
     "评测：独立分析能力\n「只给你 CPU 代码，你自己搞定」", ACCENT2),
]
for i, (title, content, purpose, color) in enumerate(levels):
    x = Inches(0.5 + i * 4.2)
    add_shape(s, x, Inches(1.8), Inches(3.8), Inches(4.5), BG_MID)
    add_shape(s, x, Inches(1.8), Inches(3.8), Inches(0.06), color)
    add_text(s, x + Inches(0.2), Inches(2.1), Inches(3.4), Inches(0.5),
             title, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.8), Inches(3.4), Inches(1.8),
             content, size=15, color=LIGHT)
    add_shape(s, x + Inches(0.2), Inches(4.5), Inches(3.4), Inches(0.02), GRAY)
    add_text(s, x + Inches(0.2), Inches(4.7), Inches(3.4), Inches(1.2),
             purpose, size=14, color=GRAY)

# =====================================================================
# SLIDE 13 — Multi-turn Agent
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "多轮 Agent 模式", "像人类工程师一样：写代码 -> 测试 -> 看 profiler -> 改进")

# Pipeline
steps = [
    ("Turn 0", "基础 Prompt\n-> LLM 生成代码", ACCENT),
    ("编译+测试", "nvcc 编译\n正确性验证", ACCENT2),
    ("Profiling", "nsys 分析\n找到瓶颈", ORANGE),
    ("Turn 1+", "反馈 + 代码\n-> LLM 改进", ACCENT3),
]
for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 3.3)
    add_shape(s, x, Inches(2.0), Inches(2.8), Inches(1.8), BG_MID)
    add_shape(s, x, Inches(2.0), Inches(2.8), Inches(0.05), color)
    add_text(s, x, Inches(2.2), Inches(2.8), Inches(0.4),
             title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.7), Inches(2.8), Inches(1.0),
             desc, size=15, color=LIGHT, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, x + Inches(2.8), Inches(2.6), Inches(0.5), Inches(0.5),
                 "->", size=24, color=GRAY, align=PP_ALIGN.CENTER)

# Metrics tracked
add_shape(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), BG_MID)
add_text(s, Inches(0.7), Inches(4.4), Inches(11.8), Inches(0.4),
         "每轮自动记录的指标：", size=18, color=ACCENT, bold=True)
metrics_text = (
    "  total_ms (端到端耗时)  |  kernel_time_ms (纯 kernel 耗时)  |  speedup_e2e (端到端加速比)\n"
    "  GPU 利用率 = kernel_time / total_time  |  自动生成 agent_metrics.csv + 趋势图 PNG\n\n"
    "  实验发现：多数模型在 3-5 轮后收敛，但偶尔会「灵光一现」大幅提升"
)
add_text(s, Inches(0.7), Inches(4.9), Inches(11.8), Inches(1.8),
         metrics_text, size=15, color=LIGHT)

# =====================================================================
# SLIDE 14 — Example walkthrough
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "举个例子：Bellman-Ford 最短路", "从 CPU 循环到 GPU 并行")

# CPU side
add_shape(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(3.0), BG_MID)
add_text(s, Inches(0.7), Inches(1.85), Inches(5.5), Inches(0.4),
         "CPU 参考实现 (串行)", size=18, color=ACCENT2, bold=True)
cpu_code = (
    "for (round = 0; round < V-1; round++)\n"
    "  for (u = 0; u < V; u++)\n"
    "    for each neighbor v of u:\n"
    "      if dist[u]+w < dist[v]:\n"
    "        dist[v] = dist[u] + w\n\n"
    "复杂度: O(V * E)\n"
    "V=500K, E=5M -> 非常慢"
)
add_text(s, Inches(0.7), Inches(2.4), Inches(5.5), Inches(2.3),
         cpu_code, size=14, color=LIGHT, font_name="Consolas")

# GPU side
add_shape(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(3.0), BG_MID)
add_text(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(0.4),
         "GPU 优化思路", size=18, color=ACCENT3, bold=True)
gpu_ideas = (
    "1. 每个线程处理一个节点 u\n"
    "   -> V 个线程同时松弛\n\n"
    "2. 用 atomicMin 处理竞争写入\n\n"
    "3. 提前终止：如果某轮没有更新\n"
    "   -> 减少无效迭代\n\n"
    "难点：同步开销、不规则内存访问"
)
add_text(s, Inches(7.2), Inches(2.4), Inches(5.5), Inches(2.3),
         gpu_ideas, size=14, color=LIGHT)

# What LLM needs to produce
add_shape(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), BG_MID)
add_text(s, Inches(0.7), Inches(5.3), Inches(11.8), Inches(0.4),
         "LLM 只需要实现：", size=18, color=ACCENT, bold=True)
tf = add_text(s, Inches(0.7), Inches(5.8), Inches(11.8), Inches(1.0),
              "solution_init(data, size)  -> 拷贝图到 GPU 显存，分配 dist[] 数组", size=15, color=LIGHT, font_name="Consolas")
add_para(tf, "solution_compute(request)  -> 设起点 dist[src]=0，跑并行 Bellman-Ford kernel", size=15, color=LIGHT, font_name="Consolas")
add_para(tf, "框架负责：读 input.bin、传参数、计时、验证结果、写 timing.json", size=14, color=GRAY)

# =====================================================================
# SLIDE 15 — Evaluation metrics
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "评测指标体系", "不只看对不对，更看快不快")

metrics = [
    ("编译通过率", "LLM 生成的代码能否被 nvcc 编译\n基础门槛，不通过则后续无法评测", ACCENT),
    ("正确率", "在 small / medium / large 三个规模上\n输出与 CPU 参考答案一致（支持容差）", ACCENT3),
    ("端到端加速比", "CPU 时间 / GPU 总时间\n包含 init + compute 全过程", ACCENT2),
    ("纯 Kernel 加速比", "CPU 时间 / GPU kernel 时间\n通过 nsys profiler 精确测量", ORANGE),
    ("GPU 利用率", "kernel 时间 / 端到端时间\n反映数据传输开销占比", RGBColor(0xBB, 0x86, 0xFC)),
]
for i, (name, desc, color) in enumerate(metrics):
    y = Inches(1.7 + i * 1.05)
    add_shape(s, Inches(0.5), y, Inches(12.3), Inches(0.9), BG_MID)
    add_shape(s, Inches(0.5), y, Inches(0.06), Inches(0.9), color)
    add_text(s, Inches(0.8), y + Inches(0.05), Inches(3.0), Inches(0.4),
             name, size=18, color=color, bold=True)
    add_text(s, Inches(3.8), y + Inches(0.1), Inches(8.5), Inches(0.7),
             desc, size=14, color=LIGHT)

# =====================================================================
# SLIDE 16 — Initial findings
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "初步发现与洞察", "（基于早期实验）")

items = [
    ("编译通过率 > 90%，但正确率参差不齐", [
        "大部分模型能写出语法正确的 CUDA，但数值精度、边界条件常出错",
        "Large 规模数据容易暴露隐藏 bug（小规模侥幸通过）",
    ]),
    ("加速比差异巨大：1x ~ 100x+", [
        "简单任务（Black-Scholes, 向量加法）：所有模型都能获得显著加速",
        "复杂任务（TSP, 收益管理 DP）：多数模型生成的代码甚至比 CPU 还慢",
    ]),
    ("L1 vs L3 差距明显", [
        "有详细优化提示时，模型表现明显更好",
        "说明当前 LLM 更擅长「执行优化方案」而非「发现优化方案」",
    ]),
    ("多轮 Agent 有效但有天花板", [
        "前 3 轮改进最大，之后趋于收敛",
        "偶尔出现「退化」—— 改着改着反而变差了",
    ]),
]
add_bullet_slide(s, items, size=17, color=LIGHT)

# =====================================================================
# SLIDE 17 — Summary & Next Steps
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
title_bar(s, "总结与下一步")

# Summary
add_shape(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5), BG_MID)
add_text(s, Inches(0.7), Inches(1.9), Inches(5.4), Inches(0.5),
         "ORBench 贡献", size=22, color=ACCENT, bold=True)
tf = add_text(s, Inches(0.7), Inches(2.5), Inches(5.4), Inches(3.5),
              "  40 道跨领域 GPU 加速任务", size=17, color=LIGHT)
add_para(tf, "  三层架构确保评测公平性", size=17, color=LIGHT)
add_para(tf, "  Prompt 三级分级评测不同能力维度", size=17, color=LIGHT)
add_para(tf, "  多轮 Agent 模式模拟真实开发流程", size=17, color=LIGHT)
add_para(tf, "  端到端自动化：一条命令跑完全流程", size=17, color=LIGHT)
add_para(tf, "", size=10, color=LIGHT)
add_para(tf, "  核心发现：LLM 能写对 CUDA，", size=17, color=ORANGE)
add_para(tf, "  但写快还需要更好的引导和迭代", size=17, color=ORANGE)

# Next steps
add_shape(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.5), BG_MID)
add_text(s, Inches(7.2), Inches(1.9), Inches(5.4), Inches(0.5),
         "下一步计划", size=22, color=ACCENT3, bold=True)
tf = add_text(s, Inches(7.2), Inches(2.5), Inches(5.4), Inches(3.5),
              "  完成全部模型 x 任务 x 级别评测", size=17, color=LIGHT)
add_para(tf, "  增加更多任务（目标 50+）", size=17, color=LIGHT)
add_para(tf, "  深入分析优化模式与失败模式", size=17, color=LIGHT)
add_para(tf, "  探索更智能的 Agent 反馈策略", size=17, color=LIGHT)
add_para(tf, "  支持多 GPU / 分布式场景", size=17, color=LIGHT)
add_para(tf, "", size=10, color=LIGHT)
add_para(tf, "  开源 Benchmark，邀请社区贡献", size=17, color=ACCENT3)

# =====================================================================
# SLIDE 18 — Thank you
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_shape(s, 0, Inches(2.5), W, Inches(2.8), BG_MID)
add_text(s, Inches(1), Inches(2.8), Inches(11), Inches(1.0),
         "Thanks!", size=54, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Q & A", size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
         "github.com/Echoscd/ORbench", size=16, color=GRAY, align=PP_ALIGN.CENTER,
         font_name="Consolas")

# ── save ────────────────────────────────────────────────────────────
out = "/root/chendong/ORbench/ORBench/ORBench_GroupMeeting.pptx"
prs.save(out)
print(f"Saved to {out}")
print(f"Total slides: {len(prs.slides)}")
